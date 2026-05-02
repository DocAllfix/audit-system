"""
Test V2 — text_handlers (Limitazione 1).

Coperture:
- Ogni handler estrae testo da file generato dinamicamente
- File mancante / corrotto → text="" + method failed_*
- Sanitize rimuove caratteri di controllo
- Cap 200k chars rispettato
- Dispatcher: docx vs doc legacy via filename suffix
- Dispatcher: xlsx vs xls
- Dispatcher: image/heic deferred to OCR
- Categoria sconosciuta → unsupported_format
"""
from __future__ import annotations

import io
from pathlib import Path

import pytest

from v2 import text_handlers as th


# ──────────────────────────────────────────────────────────────────────────────
# Sanitize
# ──────────────────────────────────────────────────────────────────────────────

def test_sanitize_removes_control_chars():
    s = "ok\x00line\x01\x02more"
    cleaned = th._sanitize(s)
    assert "\x00" not in cleaned
    assert "\x01" not in cleaned


def test_sanitize_caps_at_max():
    s = "X" * (th.MAX_OUTPUT_CHARS + 5000)
    cleaned = th._sanitize(s)
    assert len(cleaned) <= th.MAX_OUTPUT_CHARS + 50  # margine per il tag


def test_sanitize_normalizes_newlines():
    s = "a\r\nb\rc"
    assert th._sanitize(s) == "a\nb\nc"


# ──────────────────────────────────────────────────────────────────────────────
# DOCX
# ──────────────────────────────────────────────────────────────────────────────

def test_extract_docx(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("Riga uno")
    doc.add_paragraph("Riga due con accenti àèìòù")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Cella1"
    table.rows[0].cells[1].text = "Cella2"
    p = tmp_path / "test.docx"
    doc.save(str(p))

    text, method = th.extract_docx(str(p))
    assert method == th.METHOD_DOCX_OK
    assert "Riga uno" in text
    assert "Riga due con accenti" in text
    assert "Cella1" in text
    assert "àèìòù" in text


def test_extract_docx_corrupted(tmp_path):
    p = tmp_path / "bad.docx"
    p.write_bytes(b"not a docx")
    text, method = th.extract_docx(str(p))
    assert method == th.METHOD_FAILED_GENERIC
    assert text == ""


# ──────────────────────────────────────────────────────────────────────────────
# XLSX
# ──────────────────────────────────────────────────────────────────────────────

def test_extract_xlsx(tmp_path):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Bilancio"
    ws["A1"] = "Voce"
    ws["B1"] = "Valore"
    ws["A2"] = "Ricavi"
    ws["B2"] = 1500000
    ws["A3"] = "Costi"
    ws["B3"] = 980000
    p = tmp_path / "bilancio.xlsx"
    wb.save(str(p))

    text, method = th.extract_xlsx(str(p))
    assert method == th.METHOD_XLSX_OK
    assert "Bilancio" in text
    assert "Ricavi" in text
    assert "1500000" in text


def test_extract_xlsx_corrupted(tmp_path):
    p = tmp_path / "bad.xlsx"
    p.write_bytes(b"not xlsx")
    text, method = th.extract_xlsx(str(p))
    assert method == th.METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# TXT
# ──────────────────────────────────────────────────────────────────────────────

def test_extract_txt_utf8(tmp_path):
    p = tmp_path / "test.txt"
    p.write_text("Contenuto italiano àèìòù", encoding="utf-8")
    text, method = th.extract_txt(str(p))
    assert method == th.METHOD_TXT_OK
    assert "àèìòù" in text


def test_extract_txt_latin1_fallback(tmp_path):
    p = tmp_path / "latin1.txt"
    p.write_bytes("Caratteri \xe0\xe8\xec\xf2\xf9".encode("latin-1"))
    text, method = th.extract_txt(str(p))
    assert method == th.METHOD_TXT_OK
    # Caratteri latini decodificati
    assert "Caratteri" in text


# ──────────────────────────────────────────────────────────────────────────────
# RTF
# ──────────────────────────────────────────────────────────────────────────────

def test_extract_rtf(tmp_path):
    rtf_content = r"{\rtf1\ansi\ansicpg1252\cocoartf2580 Hello world from RTF}"
    p = tmp_path / "test.rtf"
    p.write_text(rtf_content, encoding="utf-8")
    text, method = th.extract_rtf(str(p))
    assert method == th.METHOD_RTF_OK
    assert "Hello world" in text


# ──────────────────────────────────────────────────────────────────────────────
# HTML
# ──────────────────────────────────────────────────────────────────────────────

def test_extract_html(tmp_path):
    html_content = """
    <html>
        <head><title>Test</title></head>
        <body>
            <h1>Documento HTML</h1>
            <p>Paragrafo di prova</p>
            <script>alert('skip')</script>
        </body>
    </html>
    """
    p = tmp_path / "test.html"
    p.write_text(html_content, encoding="utf-8")
    text, method = th.extract_html(str(p))
    assert method == th.METHOD_HTML_OK
    assert "Documento HTML" in text
    assert "Paragrafo di prova" in text
    assert "alert" not in text  # script rimosso


# ──────────────────────────────────────────────────────────────────────────────
# XML
# ──────────────────────────────────────────────────────────────────────────────

def test_extract_xml(tmp_path):
    xml_content = """<?xml version="1.0" encoding="UTF-8"?>
    <fattura>
        <azienda>
            <nome>Demo SRL</nome>
            <piva>12345678901</piva>
        </azienda>
        <importo>1500.00</importo>
    </fattura>
    """
    p = tmp_path / "test.xml"
    p.write_text(xml_content, encoding="utf-8")
    text, method = th.extract_xml(str(p))
    assert method == th.METHOD_XML_OK
    assert "Demo SRL" in text
    assert "12345678901" in text
    assert "1500.00" in text


# ──────────────────────────────────────────────────────────────────────────────
# EML
# ──────────────────────────────────────────────────────────────────────────────

def test_extract_eml(tmp_path):
    eml_content = """From: sender@example.com
To: receiver@example.com
Subject: Test Email
Date: Mon, 01 May 2026 12:00:00 +0200
Content-Type: text/plain

Corpo dell'email di test.
Riga seconda.
"""
    p = tmp_path / "test.eml"
    p.write_text(eml_content, encoding="utf-8")
    text, method = th.extract_eml(str(p))
    assert method == th.METHOD_EML_OK
    assert "sender@example.com" in text
    assert "Test Email" in text
    assert "Corpo dell'email" in text


# ──────────────────────────────────────────────────────────────────────────────
# PPTX
# ──────────────────────────────────────────────────────────────────────────────

def test_extract_pptx(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    if title:
        title.text = "Titolo Slide"

    # Aggiungi textbox manuale
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(2))
    tf = txBox.text_frame
    tf.text = "Contenuto della slide"

    p = tmp_path / "test.pptx"
    prs.save(str(p))

    text, method = th.extract_pptx(str(p))
    assert method == th.METHOD_PPTX_OK
    assert "Slide" in text
    assert "Contenuto della slide" in text


# ──────────────────────────────────────────────────────────────────────────────
# Dispatcher
# ──────────────────────────────────────────────────────────────────────────────

def test_dispatcher_docx_via_suffix(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("Dispatch test")
    p = tmp_path / "doc.docx"
    doc.save(str(p))

    file_info = {"path": str(p), "filename": "doc.docx", "category": "word"}
    text, method = th.extract_text_for_category(file_info)
    assert method == th.METHOD_DOCX_OK
    assert "Dispatch test" in text


def test_dispatcher_xlsx_via_suffix(tmp_path):
    from openpyxl import Workbook
    wb = Workbook()
    wb.active["A1"] = "OK"
    p = tmp_path / "data.xlsx"
    wb.save(str(p))

    file_info = {"path": str(p), "filename": "data.xlsx", "category": "excel"}
    text, method = th.extract_text_for_category(file_info)
    assert method == th.METHOD_XLSX_OK


def test_dispatcher_image_deferred_to_ocr(tmp_path):
    p = tmp_path / "img.png"
    p.write_bytes(b"\x89PNG\r\n\x1a\n")
    file_info = {"path": str(p), "filename": "img.png", "category": "image"}
    text, method = th.extract_text_for_category(file_info)
    assert method == "deferred_to_ocr"
    assert text == ""


def test_dispatcher_unknown_category(tmp_path):
    p = tmp_path / "x.xyz"
    p.write_bytes(b"data")
    file_info = {"path": str(p), "filename": "x.xyz", "category": "other"}
    text, method = th.extract_text_for_category(file_info)
    assert method == th.METHOD_UNSUPPORTED_FORMAT


def test_dispatcher_missing_file(tmp_path):
    file_info = {"path": "/non/esiste.docx", "filename": "x.docx", "category": "word"}
    text, method = th.extract_text_for_category(file_info)
    assert method == th.METHOD_FAILED_NO_HANDLER


def test_dispatcher_invalid_input():
    text, method = th.extract_text_for_category("not a dict")  # type: ignore
    assert method == th.METHOD_FAILED_NO_HANDLER

    text, method = th.extract_text_for_category({})
    assert method == th.METHOD_FAILED_NO_HANDLER


def test_dispatcher_txt_category(tmp_path):
    p = tmp_path / "test.txt"
    p.write_text("Plain text content", encoding="utf-8")
    file_info = {"path": str(p), "filename": "test.txt", "category": "text"}
    text, method = th.extract_text_for_category(file_info)
    assert method == th.METHOD_TXT_OK
    assert "Plain text content" in text


# ──────────────────────────────────────────────────────────────────────────────
# Mai eccezione
# ──────────────────────────────────────────────────────────────────────────────

def test_handlers_never_raise_on_garbage(tmp_path):
    """Tutti gli handler ricevono bytes random: nessuno solleva."""
    for ext in (".docx", ".xlsx", ".rtf", ".html", ".xml", ".eml", ".odt", ".pptx"):
        p = tmp_path / f"garbage{ext}"
        p.write_bytes(b"\x00\x01\x02 random garbage \xff\xfe")
        file_info = {
            "path": str(p),
            "filename": f"garbage{ext}",
            "category": {
                ".docx": "word", ".xlsx": "excel", ".rtf": "rtf",
                ".html": "html", ".xml": "xml", ".eml": "eml",
                ".odt": "odt", ".pptx": "pptx",
            }[ext],
        }
        # Non deve mai sollevare
        text, method = th.extract_text_for_category(file_info)
        # Method può essere anche OK se per caso è interpretabile, l'importante
        # è non sollevare
        assert isinstance(text, str)
        assert isinstance(method, str)
