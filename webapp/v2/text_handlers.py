"""
V2 — Text Handlers per formati non-PDF (Limitazione 1).

Replica la logica V1 di estrazione testo da docx, xlsx, txt, eml, msg, html,
xml, rtf, p7m, odt, ods, pptx — senza importare da `webapp/modules/`.

Caratteristiche:
- Mai eccezione: ogni handler ritorna `(text, method)` con method che dichiara
  l'esito (`docx_ok`, `failed_*`, ecc.)
- Cap output 200k chars (anti-blob)
- Sanitizzazione caratteri di controllo
- Dispatcher unico `extract_text_for_category(file_info)` che instrada al
  giusto handler in base alla `category` di file_triage

API pubblica:
    extract_text_for_category(file_info: dict) -> Tuple[str, str]
"""
from __future__ import annotations

import os
import re
from typing import Any, Dict, Optional, Tuple

# Cap output: stesso budget del PDF native extractor
MAX_OUTPUT_CHARS = 200_000

# Esiti standard
METHOD_DOCX_OK = "docx_ok"
METHOD_DOC_LEGACY_OK = "doc_legacy_ok"
METHOD_XLSX_OK = "xlsx_ok"
METHOD_XLS_OK = "xls_ok"
METHOD_TXT_OK = "txt_ok"
METHOD_RTF_OK = "rtf_ok"
METHOD_HTML_OK = "html_ok"
METHOD_XML_OK = "xml_ok"
METHOD_EML_OK = "eml_ok"
METHOD_MSG_OK = "msg_ok"
METHOD_ODT_OK = "odt_ok"
METHOD_ODS_OK = "ods_ok"
METHOD_PPTX_OK = "pptx_ok"
METHOD_P7M_OK = "p7m_ok"
METHOD_FAILED_NO_HANDLER = "failed_no_handler"
METHOD_FAILED_GENERIC = "failed_generic"
METHOD_FAILED_LIB_MISSING = "failed_lib_missing"
METHOD_UNSUPPORTED_FORMAT = "unsupported_format"


# ──────────────────────────────────────────────────────────────────────────────
# Sanitization
# ──────────────────────────────────────────────────────────────────────────────

def _sanitize(text: str) -> str:
    """Rimuove caratteri di controllo e tronca al cap."""
    if not text:
        return ""
    text = text.replace("\x00", "")
    text = "".join(c for c in text if ord(c) >= 32 or c in "\n\t\r")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    if len(text) > MAX_OUTPUT_CHARS:
        text = text[:MAX_OUTPUT_CHARS] + "\n[TESTO TRONCATO V2]"
    return text


# ──────────────────────────────────────────────────────────────────────────────
# DOCX (python-docx)
# ──────────────────────────────────────────────────────────────────────────────

def extract_docx(path: str) -> Tuple[str, str]:
    try:
        from docx import Document
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        doc = Document(path)
        parts = []
        # Paragrafi
        for para in doc.paragraphs:
            if para.text:
                parts.append(para.text)
        # Tabelle
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
        return _sanitize(text), METHOD_DOCX_OK
    except Exception as e:
        print(f"[V2 DOCX] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# DOC legacy (.doc) — best effort via olefile (estrazione testo grezzo)
# ──────────────────────────────────────────────────────────────────────────────

def extract_doc_legacy(path: str) -> Tuple[str, str]:
    """
    .doc è formato OLE binario. Senza Word installato si può fare solo best
    effort: estraggo lo stream "WordDocument" via olefile e tento decode.
    Su Windows con pywin32 si potrebbe usare COM, ma è fragile.
    """
    try:
        import olefile
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        if not olefile.isOleFile(path):
            return "", METHOD_FAILED_GENERIC
        ole = olefile.OleFileIO(path)
        try:
            if not ole.exists("WordDocument"):
                return "", METHOD_FAILED_GENERIC
            stream = ole.openstream("WordDocument")
            raw = stream.read()
            stream.close()
            # Filtra solo bytes ASCII printable + accentati Latin-1
            text = raw.decode("latin-1", errors="ignore")
            # Rimuovi sequenze di controllo lasciando il testo leggibile
            text = re.sub(r"[\x00-\x08\x0b-\x0c\x0e-\x1f]+", " ", text)
            text = re.sub(r"\s{3,}", "\n", text)
            return _sanitize(text), METHOD_DOC_LEGACY_OK
        finally:
            ole.close()
    except Exception as e:
        print(f"[V2 DOC] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# XLSX / XLSM (openpyxl)
# ──────────────────────────────────────────────────────────────────────────────

def extract_xlsx(path: str) -> Tuple[str, str]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        wb = load_workbook(path, data_only=True, read_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            parts.append(f"=== {sheet} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append(" | ".join(cells))
        wb.close()
        return _sanitize("\n".join(parts)), METHOD_XLSX_OK
    except Exception as e:
        print(f"[V2 XLSX] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# XLS legacy (xlrd)
# ──────────────────────────────────────────────────────────────────────────────

def extract_xls(path: str) -> Tuple[str, str]:
    try:
        import xlrd
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        book = xlrd.open_workbook(path)
        parts = []
        for sheet in book.sheets():
            parts.append(f"=== {sheet.name} ===")
            for r in range(sheet.nrows):
                row = sheet.row_values(r)
                cells = [str(c) for c in row if str(c).strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return _sanitize("\n".join(parts)), METHOD_XLS_OK
    except Exception as e:
        print(f"[V2 XLS] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# TXT / CSV
# ──────────────────────────────────────────────────────────────────────────────

def extract_txt(path: str) -> Tuple[str, str]:
    try:
        # Tenta UTF-8, fallback Latin-1
        try:
            with open(path, "r", encoding="utf-8", errors="strict") as f:
                text = f.read()
        except UnicodeDecodeError:
            with open(path, "r", encoding="latin-1", errors="ignore") as f:
                text = f.read()
        return _sanitize(text), METHOD_TXT_OK
    except Exception as e:
        print(f"[V2 TXT] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# RTF (striprtf)
# ──────────────────────────────────────────────────────────────────────────────

def extract_rtf(path: str) -> Tuple[str, str]:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read()
        text = rtf_to_text(raw)
        return _sanitize(text), METHOD_RTF_OK
    except Exception as e:
        print(f"[V2 RTF] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# HTML / HTM (BeautifulSoup)
# ──────────────────────────────────────────────────────────────────────────────

def extract_html(path: str) -> Tuple[str, str]:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            html = f.read()
        soup = BeautifulSoup(html, "lxml")
        # Rimuovi script/style
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return _sanitize(text), METHOD_HTML_OK
    except Exception as e:
        print(f"[V2 HTML] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# XML
# ──────────────────────────────────────────────────────────────────────────────

def extract_xml(path: str) -> Tuple[str, str]:
    try:
        from lxml import etree
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        tree = etree.parse(path)
        root = tree.getroot()
        # Estraggo solo il testo (no tag)
        parts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        return _sanitize("\n".join(parts)), METHOD_XML_OK
    except Exception as e:
        print(f"[V2 XML] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# EML (stdlib email)
# ──────────────────────────────────────────────────────────────────────────────

def extract_eml(path: str) -> Tuple[str, str]:
    try:
        import email
        from email import policy
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        with open(path, "rb") as f:
            msg = email.message_from_bytes(f.read(), policy=policy.default)
        parts = []
        # Header
        for header in ("From", "To", "Subject", "Date"):
            v = msg.get(header)
            if v:
                parts.append(f"{header}: {v}")
        parts.append("")
        # Body
        body = msg.get_body(preferencelist=("plain", "html"))
        if body:
            content = body.get_content()
            if isinstance(content, str):
                # Se è HTML, stripping rapido
                if body.get_content_type() == "text/html":
                    try:
                        from bs4 import BeautifulSoup
                        content = BeautifulSoup(content, "lxml").get_text(separator="\n")
                    except ImportError:
                        pass
                parts.append(content)
        return _sanitize("\n".join(parts)), METHOD_EML_OK
    except Exception as e:
        print(f"[V2 EML] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# MSG (extract_msg)
# ──────────────────────────────────────────────────────────────────────────────

def extract_msg(path: str) -> Tuple[str, str]:
    try:
        import extract_msg as em
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        msg = em.Message(path)
        parts = []
        if msg.sender:
            parts.append(f"From: {msg.sender}")
        if msg.to:
            parts.append(f"To: {msg.to}")
        if msg.subject:
            parts.append(f"Subject: {msg.subject}")
        if msg.date:
            parts.append(f"Date: {msg.date}")
        parts.append("")
        if msg.body:
            parts.append(msg.body)
        msg.close()
        return _sanitize("\n".join(parts)), METHOD_MSG_OK
    except Exception as e:
        print(f"[V2 MSG] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# ODT / ODS (odfpy)
# ──────────────────────────────────────────────────────────────────────────────

def extract_odt(path: str) -> Tuple[str, str]:
    try:
        from odf import teletype
        from odf.opendocument import load
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        doc = load(path)
        text = teletype.extractText(doc.text)
        return _sanitize(text), METHOD_ODT_OK
    except Exception as e:
        print(f"[V2 ODT] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


def extract_ods(path: str) -> Tuple[str, str]:
    try:
        from odf import opendocument, table, text as odftext, teletype
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        doc = opendocument.load(path)
        parts = []
        for sheet in doc.spreadsheet.getElementsByType(table.Table):
            sheet_name = sheet.getAttribute("name")
            parts.append(f"=== {sheet_name} ===")
            for row in sheet.getElementsByType(table.TableRow):
                cells = []
                for cell in row.getElementsByType(table.TableCell):
                    cell_text = teletype.extractText(cell).strip()
                    if cell_text:
                        cells.append(cell_text)
                if cells:
                    parts.append(" | ".join(cells))
        return _sanitize("\n".join(parts)), METHOD_ODS_OK
    except Exception as e:
        print(f"[V2 ODS] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# PPTX (python-pptx)
# ──────────────────────────────────────────────────────────────────────────────

def extract_pptx(path: str) -> Tuple[str, str]:
    try:
        from pptx import Presentation
    except ImportError:
        return "", METHOD_FAILED_LIB_MISSING
    try:
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            parts.append(line)
        return _sanitize("\n".join(parts)), METHOD_PPTX_OK
    except Exception as e:
        print(f"[V2 PPTX] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# P7M (firmato digitalmente, contiene PDF/DOCX dentro)
# ──────────────────────────────────────────────────────────────────────────────

def extract_p7m(path: str) -> Tuple[str, str]:
    """
    P7M è un wrapper PKCS#7 che incapsula un file (tipicamente PDF/DOCX).
    Approccio best-effort: cerca pattern di file embedded e tenta estrazione.
    """
    try:
        with open(path, "rb") as f:
            raw = f.read()

        # Cerca PDF magic header (%PDF-)
        pdf_start = raw.find(b"%PDF-")
        if pdf_start >= 0:
            pdf_end = raw.rfind(b"%%EOF")
            if pdf_end > pdf_start:
                pdf_bytes = raw[pdf_start:pdf_end + 5]
                # Salva temporaneo, processo via pdfium
                import tempfile
                with tempfile.NamedTemporaryFile(
                    suffix=".pdf", delete=False
                ) as tmp:
                    tmp.write(pdf_bytes)
                    tmp_path = tmp.name
                try:
                    from v2.text_extractor import extract_native_text
                    text, _ = extract_native_text(tmp_path)
                    return _sanitize(text), METHOD_P7M_OK
                finally:
                    try:
                        os.unlink(tmp_path)
                    except OSError:
                        pass

        # Cerca DOCX (PK\x03\x04 → ZIP header)
        zip_start = raw.find(b"PK\x03\x04")
        if zip_start >= 0:
            import tempfile
            with tempfile.NamedTemporaryFile(
                suffix=".docx", delete=False
            ) as tmp:
                tmp.write(raw[zip_start:])
                tmp_path = tmp.name
            try:
                text, method = extract_docx(tmp_path)
                if text:
                    return text, METHOD_P7M_OK
            finally:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass

        return "", METHOD_FAILED_GENERIC
    except Exception as e:
        print(f"[V2 P7M] {os.path.basename(path)}: {e}")
        return "", METHOD_FAILED_GENERIC


# ──────────────────────────────────────────────────────────────────────────────
# Dispatcher
# ──────────────────────────────────────────────────────────────────────────────

# Mappa category (da zip_extractor) → handler function
_DISPATCHER = {
    "word": None,    # Risolto via filename suffix
    "excel": None,   # Risolto via filename suffix
    "text": extract_txt,
    "rtf": extract_rtf,
    "html": extract_html,
    "xml": extract_xml,
    "email": extract_msg,
    "eml": extract_eml,
    "odt": extract_odt,
    "ods": extract_ods,
    "pptx": extract_pptx,
    "p7m": extract_p7m,
}


def extract_text_for_category(file_info: Dict[str, Any]) -> Tuple[str, str]:
    """
    Dispatcher principale: data una file_info dict (con `path`, `category`,
    `filename`), ritorna (text, method).

    Per i formati non supportati ritorna ("", METHOD_UNSUPPORTED_FORMAT)
    senza mai sollevare.
    """
    if not isinstance(file_info, dict):
        return "", METHOD_FAILED_NO_HANDLER

    path = file_info.get("path", "")
    if not path or not os.path.isfile(path):
        return "", METHOD_FAILED_NO_HANDLER

    category = file_info.get("category", "")
    filename = file_info.get("filename", "").lower()

    # Word: distinguo docx vs doc legacy via suffix
    if category == "word":
        if filename.endswith(".docx"):
            return extract_docx(path)
        else:
            return extract_doc_legacy(path)

    # Excel: distinguo xlsx vs xls
    if category == "excel":
        if filename.endswith(".xls"):
            return extract_xls(path)
        else:
            return extract_xlsx(path)

    # Altri formati con dispatch diretto
    handler = _DISPATCHER.get(category)
    if handler is not None:
        return handler(path)

    # Image, heic → delegato a OCR Vision (gestito altrove)
    if category in ("image", "heic"):
        return "", "deferred_to_ocr"

    # Formato sconosciuto / non gestito
    return "", METHOD_UNSUPPORTED_FORMAT
